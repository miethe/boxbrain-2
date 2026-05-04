from __future__ import annotations

import io
import os
import shutil

import pytest
from pptx import Presentation

from app.application.slide_renderer import LibreOfficeSlideRenderer


pytestmark = pytest.mark.skipif(
    os.getenv("BOXBRAIN_RUN_RENDER_TESTS") != "1",
    reason="Set BOXBRAIN_RUN_RENDER_TESTS=1 to run live LibreOffice rendering tests.",
)


def _pptx_bytes() -> bytes:
    buffer = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Live render smoke"
    presentation.save(buffer)
    return buffer.getvalue()


def test_libreoffice_renderer_produces_slide_asset() -> None:
    if not shutil.which("soffice") and not shutil.which("libreoffice"):
        pytest.skip("LibreOffice is not installed or not on PATH.")

    assets = LibreOfficeSlideRenderer().render_pptx(
        content=_pptx_bytes(),
        filename="live-render.pptx",
        slide_count=1,
    )

    assert len(assets) == 1
    assert assets[0].source_order_index == 1
    assert assets[0].render_content
    assert assets[0].thumbnail_content
